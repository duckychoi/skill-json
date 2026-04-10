#!/usr/bin/env python3
"""
PPTX 평가 스크립트 v3 - PowerPoint COM 렌더링 기반 시각 비교

평가 기준:
  1. 슬라이드 수 일치 (10%)
  2. 네이티브 요소 풍부도 - 슬라이드당 20개 이상 shape 필요 (15%)
  3. 이미지 미사용 - add_picture 금지 (15%)
  4. 텍스트 커버리지 - 슬라이드당 150자 이상 필요 (20%)
  5. 시각 유사도 - PowerPoint COM 렌더링 후 SSIM (40%)

사용법: python scripts/evaluate_pptx.py <output.pptx> <source.pdf>
"""
import sys, os, difflib, tempfile, math
sys.path.insert(0, 'C:/Users/user/pptx_libs')

try:
    import fitz; HAS_FITZ = True
except: HAS_FITZ = False
try:
    from pptx import Presentation; from pptx.enum.shapes import MSO_SHAPE_TYPE; HAS_PPTX = True
except: HAS_PPTX = False
try:
    from PIL import Image; HAS_PIL = True
except: HAS_PIL = False
try:
    from skimage.metrics import structural_similarity as ssim; import numpy as np; HAS_SSIM = True
except: HAS_SSIM = False


def score_slide_count(pptx_path, pdf_path):
    if not HAS_PPTX or not HAS_FITZ: return 0.0
    prs = Presentation(pptx_path)
    doc = fitz.open(pdf_path)
    out_n, exp_n = len(prs.slides), len(doc)
    doc.close()
    if exp_n == 0: return 0.0
    return max(0.0, 1.0 - abs(out_n/exp_n - 1.0))


def score_native_richness(pptx_path):
    if not HAS_PPTX: return 0.0
    prs = Presentation(pptx_path)
    counts = []
    for slide in prs.slides:
        native = sum(1 for s in slide.shapes
                     if s.shape_type != MSO_SHAPE_TYPE.PICTURE
                     and s.shape_type != MSO_SHAPE_TYPE.LINKED_PICTURE)
        counts.append(native)
    if not counts: return 0.0
    avg = sum(counts) / len(counts)
    # GT 기준: 슬라이드당 20개 이상이면 만점
    return min(1.0, avg / 20.0)


def score_no_image(pptx_path):
    if not HAS_PPTX: return 0.0
    prs = Presentation(pptx_path)
    total_shapes, images = 0, 0
    for slide in prs.slides:
        for s in slide.shapes:
            total_shapes += 1
            if s.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
                images += 1
    if total_shapes == 0: return 1.0
    return 1.0 - (images / total_shapes)


def score_text_coverage(pptx_path):
    if not HAS_PPTX: return 0.0
    prs = Presentation(pptx_path)
    slide_chars = []
    for slide in prs.slides:
        chars = 0
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    chars += len(p.text.strip())
            if s.has_table:
                for row in s.table.rows:
                    for cell in row.cells:
                        chars += len(cell.text.strip())
        slide_chars.append(chars)
    if not slide_chars: return 0.0
    avg = sum(slide_chars) / len(slide_chars)
    # GT 기준: 슬라이드당 평균 150자 이상이면 만점
    return min(1.0, avg / 150.0)


def _pdf_to_images(pdf_path, out_dir, dpi=150):
    if not HAS_FITZ: return []
    doc = fitz.open(pdf_path)
    paths = []
    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=dpi)
        p = os.path.join(out_dir, f"pdf_{i:03d}.png")
        pix.save(p)
        paths.append(p)
    doc.close()
    return paths


def _pptx_to_images(pptx_path, out_dir):
    """PowerPoint COM으로 PPTX→PNG"""
    try:
        import comtypes.client
        abs_pptx = os.path.abspath(pptx_path)
        slide_dir = os.path.join(os.path.abspath(out_dir), "slides")
        pp = comtypes.client.CreateObject("PowerPoint.Application")
        pp.Visible = 1
        pres = pp.Presentations.Open(abs_pptx, WithWindow=False)
        pres.SaveAs(slide_dir, 18)  # ppSaveAsPNG
        pres.Close()
        pp.Quit()
        if os.path.isdir(slide_dir):
            return sorted([os.path.join(slide_dir, f)
                          for f in os.listdir(slide_dir)
                          if f.lower().endswith('.png')])
    except Exception as e:
        print(f"[WARN] PowerPoint COM 실패: {e}", file=sys.stderr)
    return []


def score_visual_similarity(pptx_path, pdf_path):
    if not HAS_SSIM or not HAS_PIL:
        print("[WARN] SSIM 라이브러리 없음 - 0.5 반환", file=sys.stderr)
        return 0.5
    with tempfile.TemporaryDirectory() as tmp:
        pdf_dir = os.path.join(tmp, "pdf")
        pptx_dir = os.path.join(tmp, "pptx")
        os.makedirs(pdf_dir); os.makedirs(pptx_dir)
        pdf_imgs = _pdf_to_images(pdf_path, pdf_dir)
        pptx_imgs = _pptx_to_images(pptx_path, pptx_dir)
        if not pptx_imgs:
            print("[WARN] PPTX 렌더링 실패 - 0.3 반환", file=sys.stderr)
            return 0.3
        pairs = min(len(pdf_imgs), len(pptx_imgs))
        if pairs == 0: return 0.0
        total = 0.0
        for i in range(pairs):
            i1 = Image.open(pdf_imgs[i]).convert("L").resize((960, 720), Image.LANCZOS)
            i2 = Image.open(pptx_imgs[i]).convert("L").resize((960, 720), Image.LANCZOS)
            s = ssim(np.array(i1), np.array(i2), data_range=255)
            print(f"  SSIM slide {i+1}: {s:.3f}", file=sys.stderr)
            total += max(0.0, s)
        return total / pairs


WEIGHTS = {
    "slide_count": 0.10,
    "native_richness": 0.15,
    "no_image": 0.15,
    "text_coverage": 0.20,
    "visual_similarity": 0.40,
}

def evaluate(pptx_path, pdf_path):
    scores = {}
    print(f"[평가 v3] 출력: {pptx_path}", file=sys.stderr)
    print(f"[평가 v3] 원본: {pdf_path}", file=sys.stderr)

    scores["slide_count"] = score_slide_count(pptx_path, pdf_path)
    print(f"  슬라이드 수:     {scores['slide_count']:.3f}", file=sys.stderr)

    scores["native_richness"] = score_native_richness(pptx_path)
    print(f"  네이티브 풍부도: {scores['native_richness']:.3f} (20개 기준)", file=sys.stderr)

    scores["no_image"] = score_no_image(pptx_path)
    print(f"  이미지 미사용:   {scores['no_image']:.3f}", file=sys.stderr)

    scores["text_coverage"] = score_text_coverage(pptx_path)
    print(f"  텍스트 커버리지: {scores['text_coverage']:.3f} (150자/슬라이드 기준)", file=sys.stderr)

    scores["visual_similarity"] = score_visual_similarity(pptx_path, pdf_path)
    print(f"  시각 유사도:     {scores['visual_similarity']:.3f}", file=sys.stderr)

    total = sum(WEIGHTS[k] * scores[k] for k in WEIGHTS)
    print(f"  종합:            {total:.3f}", file=sys.stderr)
    return total

def main():
    if len(sys.argv) < 3:
        print("사용법: python evaluate_pptx.py <output.pptx> <source.pdf>", file=sys.stderr)
        sys.exit(1)
    pptx_path, pdf_path = sys.argv[1], sys.argv[2]
    for p in [pptx_path, pdf_path]:
        if not os.path.exists(p):
            print(f"[ERROR] 파일 없음: {p}", file=sys.stderr)
            print("0.000"); sys.exit(1)
    score = evaluate(pptx_path, pdf_path)
    print(f"{score:.6f}")

if __name__ == "__main__":
    main()
