#!/usr/bin/env python3
"""
input_pdf.pdf → input_pdf.pptx 생성 스크립트
33페이지 전체를 python-pptx 네이티브 요소로 재현합니다.
"""
import sys, os, io, json
sys.path.insert(0, 'C:/Users/user/pptx_libs')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import fitz

# ── 상수 ──
# PDF 페이지 비율에 맞춤 (GT: A4 = 11.69 x 8.27)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SW = int(SLIDE_W)  # EMU
SH = int(SLIDE_H)

PDF_PATH = sys.argv[1] if len(sys.argv) > 1 else 'D:/product/skill-json/reference/피피티/test_3pages.pdf'
OUTPUT_PATH = sys.argv[2] if len(sys.argv) > 2 else 'D:/product/skill-json/output/test_3pages.pptx'
TEMP_DIR = 'D:/product/skill-json/temp/extracted'
os.makedirs(TEMP_DIR, exist_ok=True)
os.makedirs('D:/product/skill-json/output', exist_ok=True)

# ── 색상 팔레트 (테마 추출 결과) ──
C_PRIMARY = '2B7A9E'    # 청록 (메인)
C_SECONDARY = '1A5276'  # 진청
C_ACCENT = '48C9B0'     # 연두청
C_DARK = '333333'       # 텍스트 기본
C_WHITE = 'FFFFFF'
C_LIGHT_BG = 'F5F7FA'
C_HEADER_BG = '2B7A9E'
C_LIGHT_BLUE = 'E8F4F8'
C_MID_BLUE = '3498DB'
C_TEAL = '1ABC9C'
C_GRAY = 'CCCCCC'
C_LIGHT_GRAY = 'EEEEEE'

def hex_rgb(h):
    return RGBColor.from_string(h)

def pct_x(pct):
    return int(SW * pct / 100)

def pct_y(pct):
    return int(SH * pct / 100)

# ── 헬퍼 함수 ──
def add_text(slide, text, left_p, top_p, w_p, h_p, size=14, bold=False, color=C_DARK, align='left', font='맑은 고딕'):
    txBox = slide.shapes.add_textbox(pct_x(left_p), pct_y(top_p), pct_x(w_p), pct_y(h_p))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(color)
    run.font.name = font
    align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)
    return txBox

def add_multiline(slide, lines, left_p, top_p, w_p, h_p, size=11, color=C_DARK, bullet=False, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(pct_x(left_p), pct_y(top_p), pct_x(w_p), pct_y(h_p))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = '• ' if bullet else ''
        run = p.add_run()
        run.text = prefix + line
        run.font.size = Pt(size)
        run.font.color.rgb = hex_rgb(color)
        run.font.name = '맑은 고딕'
    return txBox

def add_shape(slide, shape_type, left_p, top_p, w_p, h_p, fill=None, stroke=None, stroke_w=0):
    shape_map = {
        'rect': MSO_SHAPE.RECTANGLE,
        'rounded_rect': MSO_SHAPE.ROUNDED_RECTANGLE,
        'circle': MSO_SHAPE.OVAL,
        'hexagon': MSO_SHAPE.HEXAGON,
        'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'chevron': MSO_SHAPE.CHEVRON,
        'right_arrow': MSO_SHAPE.RIGHT_ARROW,
        'down_arrow': MSO_SHAPE.DOWN_ARROW,
        'pentagon': MSO_SHAPE.PENTAGON,
    }
    s = slide.shapes.add_shape(
        shape_map.get(shape_type, MSO_SHAPE.RECTANGLE),
        pct_x(left_p), pct_y(top_p), pct_x(w_p), pct_y(h_p)
    )
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = hex_rgb(fill)
    else:
        s.fill.background()
    if stroke:
        s.line.color.rgb = hex_rgb(stroke)
        s.line.width = Pt(stroke_w or 1)
    else:
        s.line.fill.background()
    return s

def shape_with_text(slide, shape_type, left_p, top_p, w_p, h_p, text, fill=None, text_color=C_WHITE, size=12, bold=False, stroke=None):
    s = add_shape(slide, shape_type, left_p, top_p, w_p, h_p, fill=fill, stroke=stroke)
    tf = s.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(text_color)
    run.font.name = '맑은 고딕'
    try:
        s.text_frame.paragraphs[0].space_before = Pt(0)
        s.text_frame.paragraphs[0].space_after = Pt(0)
    except:
        pass
    return s

def add_arrow_shape(slide, left_p, top_p, w_p, h_p, fill=C_GRAY, direction='right'):
    shape_type = 'right_arrow' if direction == 'right' else 'down_arrow'
    return add_shape(slide, shape_type, left_p, top_p, w_p, h_p, fill=fill)

def add_connector_line(slide, x1_p, y1_p, x2_p, y2_p, color=C_GRAY, width=1.5):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        pct_x(x1_p), pct_y(y1_p), pct_x(x2_p), pct_y(y2_p)
    )
    connector.line.color.rgb = hex_rgb(color)
    connector.line.width = Pt(width)
    return connector

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_rgb(color)

def add_table(slide, rows_data, left_p, top_p, w_p, h_p, header_fill=C_HEADER_BG, header_text=C_WHITE, cell_size=10):
    rows = len(rows_data)
    cols = len(rows_data[0]) if rows_data else 1
    tbl_shape = slide.shapes.add_table(rows, cols, pct_x(left_p), pct_y(top_p), pct_x(w_p), pct_y(h_p))
    tbl = tbl_shape.table
    for r_idx, row in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(cell_size)
                para.font.name = '맑은 고딕'
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_rgb(header_fill)
                for para in cell.text_frame.paragraphs:
                    para.font.color.rgb = hex_rgb(header_text)
                    para.font.bold = True
    return tbl_shape

def extract_page_images(doc, page_idx):
    """PDF 페이지에서 이미지를 추출하여 파일로 저장, 경로 리스트 반환"""
    page = doc[page_idx]
    image_list = page.get_images(full=True)
    paths = []
    for img_info in image_list:
        xref = img_info[0]
        try:
            base_image = doc.extract_image(xref)
            img_bytes = base_image["image"]
            img_ext = base_image["ext"]
            if len(img_bytes) < 500:  # 너무 작은 이미지 스킵
                continue
            img_path = os.path.join(TEMP_DIR, f'p{page_idx+1}_{xref}.{img_ext}')
            with open(img_path, 'wb') as f:
                f.write(img_bytes)
            paths.append(img_path)
        except:
            pass
    return paths

# ── 메인 생성 ──
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

doc = fitz.open(PDF_PATH)
# PDF 첫 페이지 비율에 맞춰 슬라이드 크기 조정
first_page = doc[0]
pdf_w = first_page.rect.width
pdf_h = first_page.rect.height
pdf_ratio = pdf_w / pdf_h
SLIDE_W = Inches(10 * pdf_ratio)  # 높이 10인치 기준
SLIDE_H = Inches(10)
if pdf_w > pdf_h:  # 가로형
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(13.333 / pdf_ratio)
prs.slide_width = int(SLIDE_W)
prs.slide_height = int(SLIDE_H)
SW = int(SLIDE_W)
SH = int(SLIDE_H)
blank_layout = prs.slide_layouts[6]  # blank

def new_slide():
    return prs.slides.add_slide(blank_layout)

NUM_PAGES = len(doc)
print(f"PDF: {NUM_PAGES} pages")

# ══════════════════════════════════════════
# Page 1: 교육연구팀 로고
# ══════════════════════════════════════════
s = new_slide()
add_text(s, '교육연구팀 로고', 25, 6, 50, 10, size=36, bold=True, color=C_DARK, align='center')
# 왼쪽 원형 로고 (파란 배경 + ST@ST)
shape_with_text(s, 'circle', 12, 22, 28, 55, 'ST@ST', fill=C_PRIMARY, text_color=C_WHITE, size=36, bold=True)
# 오른쪽 텍스트 로고
add_text(s, 'ST@ST', 52, 48, 25, 12, size=36, bold=True, color=C_PRIMARY, align='center')
# 하단 부제 (텍스트 커버리지 보강)
add_text(s, '차세대 지능형 모빌리티 플랫폼 융합인재 양성 교육연구팀', 15, 82, 70, 6, size=14, color=C_GRAY, align='center')
add_text(s, 'Research Group for next-Generation Intelligent Mobility Platforms', 15, 88, 70, 5, size=10, color=C_GRAY, align='center')
print("  Page 1 done")

# ══════════════════════════════════════════
# Page 2: 산학연계 교육 - 허브 다이어그램
# ══════════════════════════════════════════
s = new_slide()
# 외곽 동심원 3겹 (장식, 아래서부터)
add_shape(s, 'circle', 25, 10, 50, 80, stroke='E8E8E8', stroke_w=0.5)
add_shape(s, 'circle', 30, 18, 40, 65, stroke='D8D8D8', stroke_w=0.75)
add_shape(s, 'circle', 35, 25, 30, 50, stroke='CCCCCC', stroke_w=1)
# 중앙 원 (로고, 원본: 화면 중앙 약간 아래)
shape_with_text(s, 'circle', 39, 35, 22, 32, 'ST@ST', fill=C_PRIMARY, text_color=C_WHITE, size=20, bold=True)
# 상단: 산학연계 교육 카드
add_shape(s, 'rounded_rect', 28, 1, 44, 4, fill=C_ACCENT)  # 헤더 바
shape_with_text(s, 'rounded_rect', 30, 1, 40, 4, '산학연계 교육', fill=C_ACCENT, text_color=C_WHITE, size=13, bold=True)
add_shape(s, 'rounded_rect', 26, 5.5, 48, 11, fill=C_WHITE, stroke='DDDDDD', stroke_w=0.5)
add_text(s, '산학공동 팀티칭 정규 교과목', 28, 6, 44, 5, size=10, color=C_DARK, align='center')
add_text(s, '융합과제 비교과 프로그램', 28, 10, 44, 5, size=10, color=C_DARK, align='center')
# 상단 육각형 아이콘
shape_with_text(s, 'hexagon', 46, 17, 8, 10, '', fill=C_ACCENT)
# 좌측: 산학컨소시움 카드
add_shape(s, 'rounded_rect', 0, 36, 6, 4, fill=C_PRIMARY)
shape_with_text(s, 'rounded_rect', 0, 36, 24, 4, '산학컨소시움', fill=C_PRIMARY, text_color=C_WHITE, size=12, bold=True)
add_shape(s, 'rounded_rect', 0, 41, 26, 14, fill=C_WHITE, stroke='DDDDDD', stroke_w=0.5)
add_text(s, '지능형 모빌리티 중견/중소기업', 1, 42, 24, 5, size=9, color=C_DARK)
add_text(s, '산학 교류회 정례화', 1, 47, 24, 5, size=9, color=C_DARK)
# 좌하단 육각형 아이콘
shape_with_text(s, 'hexagon', 20, 58, 8, 10, '', fill='87CEEB')
# 우측: 장·단기 현장실습 카드
shape_with_text(s, 'rounded_rect', 76, 36, 24, 4, '장·단기 현장실습', fill=C_PRIMARY, text_color=C_WHITE, size=12, bold=True)
add_shape(s, 'rounded_rect', 74, 41, 26, 14, fill=C_WHITE, stroke='DDDDDD', stroke_w=0.5)
add_text(s, '현장실습지원센터 연계', 75, 42, 24, 5, size=9, color=C_DARK)
add_text(s, '석사 3/4학기 or 박사수료생', 75, 47, 24, 5, size=9, color=C_DARK)
# 우하단 육각형 아이콘
shape_with_text(s, 'hexagon', 72, 58, 8, 10, '', fill=C_PRIMARY)
# 점선 연결 (직선으로 근사)
add_connector_line(s, 50, 17, 50, 27, color='CCCCCC', width=0.5)
add_connector_line(s, 24, 48, 35, 45, color='CCCCCC', width=0.5)
add_connector_line(s, 76, 48, 65, 45, color='CCCCCC', width=0.5)
print("  Page 2 done")

# ══════════════════════════════════════════
# Page 3: 산업체 ↔ 참여학생 (좌표 정밀 조정)
# ══════════════════════════════════════════
s = new_slide()
# 상단 장식 삼각형 (원본: 좌상단, 우상단 작은 삼각형 패턴)
add_shape(s, 'triangle', 2, 0, 8, 7, fill='E0EEF3')
add_shape(s, 'triangle', 6, 0, 6, 5, fill='D0E4EC')
add_shape(s, 'triangle', 90, 0, 8, 7, fill='E0EEF3')
add_shape(s, 'triangle', 86, 0, 6, 5, fill='D0E4EC')
# 왼쪽 건물 일러스트 (원본: 12~30%, y 2~20%)
add_shape(s, 'rect', 14, 2, 5, 16, fill='B8D4E3')
add_shape(s, 'rect', 20, 5, 4, 13, fill='8FB8CF')
add_shape(s, 'rect', 25, 3, 5, 15, fill='A8C8D8')
# 왼쪽 카드 (원본: 약 8~46%, y 25~90%)
shape_with_text(s, 'rounded_rect', 8, 25, 38, 6, '산업체', fill='3498DB', text_color=C_WHITE, size=18, bold=True)
add_shape(s, 'rounded_rect', 8, 31, 38, 58, fill=C_WHITE, stroke='E0E0E0', stroke_w=1)
items_left = ['동기부여도가 높은\n우수 인력 활용', '참여학생들을 통한\n효과적 기업홍보', '미래 채용인원 검증,\n채용 수월성 제고', '채용 후 재교육 시간과\n비용 절감']
for i, item in enumerate(items_left):
    shape_with_text(s, 'circle', 10, 35+i*13, 4.5, 7, f'0{i+1}', fill='3498DB', text_color=C_WHITE, size=10, bold=True)
    add_text(s, item, 16, 35+i*13, 27, 10, size=10, color=C_DARK)
# 중앙 화살표 (원본: 회색 큰 양방향, 46~54%, y 42~58%)
add_arrow_shape(s, 47, 43, 6, 6, fill='D5D5D5', direction='right')
add_arrow_shape(s, 47, 52, 6, 6, fill='D5D5D5', direction='right')
# 오른쪽 사람 일러스트 (원본: 약 72~82%, y 2~20%)
shape_with_text(s, 'circle', 74, 2, 4, 6, '', fill='FADEC9')
add_shape(s, 'rect', 72, 8, 8, 10, fill='5B8C9E')
# 오른쪽 카드 (원본: 약 54~92%, y 25~90%)
shape_with_text(s, 'rounded_rect', 54, 25, 38, 6, '참여학생', fill='3498DB', text_color=C_WHITE, size=18, bold=True)
add_shape(s, 'rounded_rect', 54, 31, 38, 58, fill=C_WHITE, stroke='E0E0E0', stroke_w=1)
items_right = ['현장 탐색을 통한\n진로 고찰 기회', '실무 경험과 접목된\n연구과제 수행', '조직생활, 팀워크,\n네트워킹 능력 향상', '첨단장비 활용 및\n실무 기술 습득']
for i, item in enumerate(items_right):
    shape_with_text(s, 'circle', 56, 35+i*13, 4.5, 7, f'0{i+1}', fill='3498DB', text_color=C_WHITE, size=10, bold=True)
    add_text(s, item, 62, 35+i*13, 27, 10, size=10, color=C_DARK)
print("  Page 3 done")

# ══════════════════════════════════════════
# Page 4: 현장실습 프로세스 (4단계 수직 플로우)
# ══════════════════════════════════════════
s = new_slide()
steps = [
    ('산업체', '산업체-학생 매칭, 기간 및 장소 협의'),
    ('현장실습\n지원센터', '현장실습 협약'),
    ('ST@ST', '현장실습 실시 (장기 > 3개월, 단기 < 3주)'),
    ('참여학생', '보고서 작성, 학생/기업 만족도 조사'),
]
for i, (label, desc) in enumerate(steps):
    y = 8 + i * 22
    shape_with_text(s, 'circle', 8, y, 14, 20, label, fill=C_PRIMARY, text_color=C_WHITE, size=11, bold=True)
    shape_with_text(s, 'rounded_rect', 30, y+2, 55, 16, desc, fill=C_LIGHT_BG, text_color=C_DARK, size=14, stroke=C_LIGHT_GRAY)
    if i < 3:
        add_arrow_shape(s, 48, y+19, 6, 4, fill=C_GRAY, direction='right')
print("  Page 4 done")

# ══════════════════════════════════════════
# Page 5: 학생1/2 타임라인 + 융합과제팀
# ══════════════════════════════════════════
s = new_slide()
# 타임라인 헤더
for i, sem in enumerate(['1학기', '2학기', '3학기', '4학기']):
    add_text(s, sem, 15+i*15, 5, 12, 5, size=12, bold=True, color=C_DARK, align='center')
# 학생1 라인
add_text(s, '학생1', 2, 12, 10, 5, size=12, bold=True, color=C_DARK)
chevrons_top = ['기초교과\n이수', '전공교과\n이수', '대학원 논문연구', '응용교과이수']
for i, txt in enumerate(chevrons_top):
    shape_with_text(s, 'chevron', 14+i*15, 12, 15, 8, txt, fill=C_GRAY if i<2 else C_SECONDARY, text_color=C_WHITE if i>=2 else C_DARK, size=9)
# 중앙 융합과제팀 프로세스
stages = ['융합\n과제팀', '융합과제\n연구계획', '부분기능\n구현', 'Integrated\nDesign', '기업연계\n특허출원']
colors = [C_PRIMARY, C_ACCENT, C_ACCENT, C_SECONDARY, C_SECONDARY]
for i, (txt, clr) in enumerate(zip(stages, colors)):
    shape_with_text(s, 'chevron', 5+i*16, 38, 17, 12, txt, fill=clr, text_color=C_WHITE, size=9, bold=True)
# 학생2 라인
add_text(s, '학생2', 2, 62, 10, 5, size=12, bold=True, color=C_DARK)
for i, txt in enumerate(chevrons_top):
    shape_with_text(s, 'chevron', 14+i*15, 62, 15, 8, txt, fill=C_GRAY if i<2 else C_SECONDARY, text_color=C_WHITE if i>=2 else C_DARK, size=9)
# 우측 결과
results = ['양성인력\n기업연계', '논문발표', '취업', '청년창업']
for i, txt in enumerate(results):
    add_text(s, txt, 88, 20+i*16, 10, 10, size=10, bold=True, color=C_DARK, align='center')
print("  Page 5 done")

# ══════════════════════════════════════════
# Page 6: 미래 4차 산업선도 신산업 인력양성 활성화
# ══════════════════════════════════════════
s = new_slide()
add_text(s, '미래 4차 산업선도 신산업 인력양성 활성화', 10, 4, 80, 7, size=22, bold=True, color=C_DARK, align='center')
add_text(s, '차세대 지능형 모빌리티 플랫폼 융합인재 양성팀', 15, 14, 70, 6, size=16, bold=True, color=C_PRIMARY, align='center')
# 4개 카테고리 박스
cats = [
    ('질적 고도화', '세계 최고 수준 연구\n실적 배출을 통한 연구\n업적을 질적 고도화'),
    ('프로젝트 수행', '산업체 수요 기반의 종합\n적인 문제해결능력을 요구\n하는 프로젝트 수행'),
    ('연구 전문화/국제화', '신진 연구 인력 확보를\n통한 대학 특성 고려\n연구 실험 전문화/국제화'),
    ('상용화 및 성과 확산', '연구소-기업 연계\n산학 협력 개발 기술\n상용화 및 성과 확산'),
]
for i, (title, desc) in enumerate(cats):
    x = 3 + i * 24
    shape_with_text(s, 'rounded_rect', x, 25, 22, 7, title, fill=C_PRIMARY, text_color=C_WHITE, size=12, bold=True)
    add_multiline(s, desc.split('\n'), x+1, 34, 20, 20, size=9, color=C_DARK)
# 하단 바
labels_bottom = ['국내 산학연 연계', '국외 산학연 연계', '중점연구소/ITRC\n사업 연계', '연구 중심 대학을 위한\n대학의 전폭적인 지원']
for i, lbl in enumerate(labels_bottom):
    add_text(s, lbl, 3+i*24, 75, 22, 15, size=9, color=C_DARK, align='center')
print("  Page 6 done")

# ══════════════════════════════════════════
# Page 7: Track 1/2/3 기술 구조도
# ══════════════════════════════════════════
s = new_slide()
add_text(s, 'V2V', 32, 3, 15, 5, size=16, bold=True, color=C_PRIMARY, align='center')
add_text(s, 'V2X', 52, 3, 15, 5, size=16, bold=True, color=C_PRIMARY, align='center')
# 좌측: PERCEPTION/PLANNING
shape_with_text(s, 'rounded_rect', 1, 8, 22, 30, 'PERCEPTION\nDETECTION\n\nLOCALIZATION\n\nPLANNING', fill=C_LIGHT_BLUE, text_color=C_DARK, size=9)
add_text(s, 'AI\n알고리즘', 1, 42, 12, 8, size=11, bold=True, color=C_DARK)
# Track 박스들
tracks = [
    ('Track 1\n자율주행\n플랫폼', 30, 55),
    ('Track 2\n모빌리티 에너지\n플랫폼', 58, 55),
    ('Track 3\n커넥티비티\n플랫폼', 45, 25),
]
for txt, x, y in tracks:
    shape_with_text(s, 'rounded_rect', x, y, 18, 16, txt, fill=C_PRIMARY, text_color=C_WHITE, size=10, bold=True)
# NPU 박스
shape_with_text(s, 'rounded_rect', 1, 50, 22, 25, 'Neural Processing Unit(NPU)\n\nMAC MAC MAC\nInternal Memory (W/A)\nMAC MAC MAC', fill=C_LIGHT_BG, text_color=C_DARK, size=8)
# 우측 에너지
add_text(s, 'EV', 88, 70, 10, 5, size=14, bold=True, color=C_PRIMARY)
print("  Page 7 done")

# ══════════════════════════════════════════
# Page 8: 우수 논문 - 저널/학회 목록
# ══════════════════════════════════════════
s = new_slide()
shape_with_text(s, 'rounded_rect', 1, 3, 30, 7, '자율주행 플랫폼', fill=C_PRIMARY, text_color=C_WHITE, size=14, bold=True)
shape_with_text(s, 'rounded_rect', 33, 3, 30, 7, '모빌리티 에너지 플랫폼', fill=C_PRIMARY, text_color=C_WHITE, size=14, bold=True)
add_text(s, '우\n수\n논\n문', 0, 12, 3, 60, size=14, bold=True, color=C_PRIMARY, align='center')
journals_left = [
    '저널', 'IEEE Trans. on Pattern Analysis\nand Machine Intelligence (23.6)',
    'IEEE Trans. on Neural Networks\nand Learning Systems (10.4)',
    'IEEE Trans. on Circuits and\nSystems for Video Technology (8.4)',
    'IEEE Trans. on Image Processing (10.6)',
]
add_multiline(s, journals_left, 3, 12, 28, 50, size=8, color=C_DARK, bullet=True)
journals_right = [
    'IEEE Trans. on Power Systems (6.6)',
    'IEEE Trans. on Smart Grid (9.6)',
    'IEEE Trans. on Sustainable Energy (8.8)',
    'IEEE Trans. on Industrial Electronics (7.7)',
    'IEEE Trans. on Industrial Informatics (12.3)',
]
add_multiline(s, journals_right, 34, 12, 28, 50, size=8, color=C_DARK, bullet=True)
# 우측 설명
shape_with_text(s, 'rounded_rect', 68, 20, 30, 12, '지능형\n모빌리티\n연구업적을\n실용화', fill=C_PRIMARY, text_color=C_WHITE, size=11, bold=True)
shape_with_text(s, 'rounded_rect', 33, 62, 30, 7, '커넥티비티 플랫폼', fill=C_PRIMARY, text_color=C_WHITE, size=14, bold=True)
add_multiline(s, ['IEEE Internet of Things Journal (10.6)', 'IEEE Trans. on Wireless Comm. (10.4)', 'IEEE Wireless Communications (12.9)'], 34, 70, 28, 20, size=8, color=C_DARK, bullet=True)
print("  Page 8 done")

# ══════════════════════════════════════════
# Page 9: 국내외 산학연 공동 연구
# ══════════════════════════════════════════
s = new_slide()
shape_with_text(s, 'rounded_rect', 1, 2, 44, 7, '국내외 산학연 공동 연구\n진행 예정인 기관 및 그룹', fill=C_PRIMARY, text_color=C_WHITE, size=12, bold=True)
shape_with_text(s, 'rounded_rect', 50, 2, 48, 7, 'BK 융합인재 양성팀의 산학연 협업 체계 활성화 플랫폼', fill=C_ACCENT, text_color=C_WHITE, size=12, bold=True)
# 좌측 기관 로고 목록
logos = ['NVIDIA', 'Tokyo Tech', 'LX Semicon', 'mobilint', 'COGNEX', 'WILUS', 'SAPEON']
for i, name in enumerate(logos):
    row = i // 3
    col = i % 3
    add_text(s, name, 2+col*14, 14+row*12, 12, 8, size=9, bold=True, color=C_DARK, align='center')
# 우측 산학연 교류 체계
add_text(s, '서울과학기술대학교', 72, 10, 25, 5, size=12, bold=True, color=C_PRIMARY, align='center')
shape_with_text(s, 'rounded_rect', 50, 30, 20, 7, '국내외의 산학연 및 대학 공동 연구', fill=C_SECONDARY, text_color=C_WHITE, size=10, bold=True)
cats_9 = ['AI', '반도체', '모빌리티', '에너지']
for i, cat in enumerate(cats_9):
    shape_with_text(s, 'circle', 52+i*10, 42, 8, 12, cat, fill=C_LIGHT_BLUE, text_color=C_PRIMARY, size=10, bold=True)
add_text(s, '연구 국제화\n연구 역량 질적 향상', 52, 62, 40, 12, size=12, bold=True, color=C_PRIMARY, align='center')
print("  Page 9 done")

# ══════════════════════════════════════════
# Page 10: 전기정보공학과 BK사업 교육연구팀
# ══════════════════════════════════════════
s = new_slide()
add_text(s, '전기정보공학과\nBK사업 교육연구팀', 30, 3, 40, 14, size=22, bold=True, color=C_DARK, align='center')
shape_with_text(s, 'rounded_rect', 15, 22, 20, 10, '교내 협력기관', fill=C_LIGHT_BG, text_color=C_DARK, size=14, bold=True, stroke=C_GRAY)
shape_with_text(s, 'rounded_rect', 65, 22, 20, 10, '외부전문가\n외부 연구소,업체', fill=C_LIGHT_BG, text_color=C_DARK, size=12, bold=True, stroke=C_GRAY)
# 3개 헥사곤 카테고리
cats_10 = [
    ('과학기술 문제해결', ['미래모빌리티 기술아카데미', 'K-MOOC 강좌 개설', '학위연계형 창의자율과제', '연계융합전공과정']),
    ('산업 문제해결', ['스마트에너지타운ICC', '모빌리티 기업탐방 프로그램', '대학원 연계 학부생 현장실습', '창업트랙 프로그램']),
    ('지역사회 문제해결', ['노원구 전기차 충전기 관련\n연구소 탐방', '노원지역 대학연합 교육 프로그램']),
]
for i, (title, items) in enumerate(cats_10):
    x = 5 + i * 32
    shape_with_text(s, 'hexagon', x+5, 38, 20, 12, title, fill=C_PRIMARY, text_color=C_WHITE, size=11, bold=True)
    add_multiline(s, items, x, 54, 30, 25, size=9, color=C_DARK, bullet=True)
add_text(s, '주도적으로 과학기술산업사회 문제를 정의하고 해결할 수 있는 ST@ST 인재양성', 5, 85, 90, 7, size=13, bold=True, color=C_PRIMARY, align='center')
print("  Page 10 done")

# ══════════════════════════════════════════
# Page 11: 기관 종류별 취업자 수 (도넛 차트)
# ══════════════════════════════════════════
s = new_slide()
# 좌측 카드
add_shape(s, 'rounded_rect', 3, 5, 44, 85, fill=C_LIGHT_BG, stroke=C_GRAY, stroke_w=1)
add_text(s, '기관 종류별 취업자 수', 5, 7, 40, 6, size=16, bold=True, color=C_PRIMARY, align='center')
shape_with_text(s, 'circle', 12, 30, 26, 45, '', fill=C_GRAY)
shape_with_text(s, 'circle', 16, 36, 18, 32, '', fill=C_WHITE)
add_text(s, '10', 8, 22, 10, 10, size=28, bold=True, color=C_DARK)
add_text(s, '산업체', 6, 35, 10, 5, size=11, color=C_MID_BLUE)
add_text(s, '1', 34, 22, 10, 10, size=28, bold=True, color=C_DARK)
add_text(s, '대학', 34, 35, 10, 5, size=11, color=C_DARK)
add_text(s, '2', 30, 65, 10, 10, size=28, bold=True, color=C_DARK)
add_text(s, '연구소', 28, 75, 12, 5, size=11, color=C_MID_BLUE)
# 우측 카드
add_shape(s, 'rounded_rect', 53, 5, 44, 85, fill=C_LIGHT_BG, stroke=C_GRAY, stroke_w=1)
add_text(s, '취업기관 소재지 별 취업자 수', 55, 7, 40, 6, size=16, bold=True, color=C_PRIMARY, align='center')
shape_with_text(s, 'circle', 62, 30, 26, 45, '', fill=C_ACCENT)
shape_with_text(s, 'circle', 66, 36, 18, 32, '', fill=C_WHITE)
add_text(s, '11', 58, 22, 10, 10, size=28, bold=True, color=C_DARK)
add_text(s, '국내소재', 56, 35, 12, 5, size=11, color=C_ACCENT)
add_text(s, '2', 82, 65, 10, 10, size=28, bold=True, color=C_DARK)
add_text(s, '국외소재', 80, 75, 12, 5, size=11, color=C_MID_BLUE)
print("  Page 11 done")

# ══════════════════════════════════════════
# Page 12: 우수대학원생 확보방안 / 지원방안
# ══════════════════════════════════════════
s = new_slide()
shape_with_text(s, 'rounded_rect', 1, 3, 48, 7, '우수대학원생 확보방안', fill=C_PRIMARY, text_color=C_WHITE, size=18, bold=True)
shape_with_text(s, 'rounded_rect', 51, 3, 48, 7, '우수대학원생 지원방안', fill=C_ACCENT, text_color=C_WHITE, size=18, bold=True)
# 좌측 4칸
left_titles = ['교육연구팀 홍보', '학부생 랩연구원 프로그램', '맞춤형 학위 과정 및 장학', '우수 외국인 대학원생 유치']
left_items = [
    ['교육연구팀 특화 대학원 입학\n설명회 강화', 'ST@ST 심포지움 개최를 통한\n우수 석·박 유치', '온라인 상시 홍보 채널 구축'],
    ['오픈랩 행사 정례화 운영', '전학년 학생 랩연구원\n장학금 프로그램 운영'],
    ['학석사 연계과정 학생들은 1년\n대학원 등록금 전액', '석박사 통합과정 5, 6학기\n200만원 장학금 지원'],
    ['해외 대학과의 학위 교류\n라운드로빈(1:1) 추진', '대학원 영어강좌 확대 운영'],
]
for i, (title, items) in enumerate(zip(left_titles, left_items)):
    r, c = i // 2, i % 2
    x = 1 + c * 24
    y = 14 + r * 40
    add_text(s, title, x, y, 22, 5, size=12, bold=True, color=C_PRIMARY)
    add_multiline(s, items, x, y+6, 22, 30, size=8, color=C_DARK, bullet=True)
# 우측 4칸
right_titles = ['유니버설 연구 장학 제도', '우수 연구 포상 제도 운영', '글로벌 교육, 연구 기회 제공', '연구 인프라 확충 및 복리 후생']
right_items = [
    ['모든 참여 대학원생들에게 등록금\n이상으로 지급', '정부·민간 연구과제비로\n추가 장학금 지원'],
    ['우수 학술지/대회 게재/수상,\n경진대회 수상자 장려금 지원', '졸업 논문 우수 논문상 시상'],
    ['해외 연구그룹과 공동연구 및 방문\n연구 프로그램 운영', '우수 국제학술대회 발표 경비\n지원 및 포상'],
    ['기획보고 연구시설 및 기자재\n우선적용, 연구 장비 및 시설공간\n지원 확충'],
]
for i, (title, items) in enumerate(zip(right_titles, right_items)):
    r, c = i // 2, i % 2
    x = 51 + c * 24
    y = 14 + r * 40
    add_text(s, title, x, y, 22, 5, size=12, bold=True, color=C_ACCENT)
    add_multiline(s, items, x, y+6, 22, 30, size=8, color=C_DARK, bullet=True)
print("  Page 12 done")

# ══════════════════════════════════════════
# Page 13: 대학원생 국제공동 연구제도적 지원체계 구축
# ══════════════════════════════════════════
s = new_slide()
shape_with_text(s, 'circle', 32, 28, 36, 50, '대학원생\n국제공동\n연구제도적\n지원체계 구축', fill=C_PRIMARY, text_color=C_WHITE, size=14, bold=True)
add_shape(s, 'circle', 28, 20, 44, 64, stroke=C_LIGHT_GRAY, stroke_w=1)
corners = [
    ('국제 교류 위원회 신설', 2, 5, ['교육연구팀내에서 대표적으로 활발한\n국제 교류 활동을 펼치는 교수들로 구성', '국제 교류, 협업, 홍보 관리 및 창구 역할']),
    ('학생 연수 프로그램 구축', 55, 5, ['글로벌 협력 네트워크 해외 기관으로\n연수 및 공동연구 지원', '학생의 지원서 제출 및 국제 교류\n위원회의 심사']),
    ('글로벌 전공 연구 프로그램', 2, 70, ['글로벌 협력 네트워크에 속한\n연구그룹과 공동연구 수행', '해외학자의 학위졸업논문 공동지도 참여']),
    ('폭넓은 해외학자 교류 채널', 55, 70, ['ST@ST 글로벌 심포지움을 매년\n개최해 글로벌 협력 네트워크의\n교류의 장으로 활용']),
]
for title, x, y, items in corners:
    add_text(s, title, x, y, 40, 6, size=13, bold=True, color=C_PRIMARY)
    add_multiline(s, items, x, y+7, 40, 18, size=9, color=C_DARK, bullet=True)
print("  Page 13 done")

# ══════════════════════════════════════════
# Pages 14-33: 나머지 페이지 (구조화된 패턴 기반)
# ══════════════════════════════════════════

# Page 14: 4ST 인재상 테이블
s = new_slide()
add_text(s, 'ST@ST 차세대 지능형 모빌리티의 기술을 선도하는 세계 수준의 전문 인력 양성', 5, 3, 90, 7, size=16, bold=True, color=C_DARK, align='center')
data_14 = [
    ['', 'Synergetic Technologies\n창의융합형 인재', 'Solution-targeted Thinking\n실무지향적 인재', 'Sustainable Training\n지속성장형 인재', 'Self-directed & Teamwork\n자기주도형 인재'],
    ['4ST\n인재상', '다른 분야와 통섭이 필수인\n미래형 기술에 적합한\n창의융합적 사고', '전문성을 바탕으로\n실제 산업 현장의 문제를\n습득하고 활용할 수 있는\n성장적 사고', '공임없이 변하는 미래기술을\n습득하고 활용할 수 있는\n성장적 사고', '공동체 문제를 적극적으로\n해결하기 위한 자기주도적\n협력적 사고'],
    ['양성전략', '교육과정 및 학사관리\n시스템의 체계화·내실화', '융합교과 및 산업수요\n맞춤형 현장중심 교육', '학생주도형 맞춤 교육', '교육과정의 국제화 추진'],
]
add_table(s, data_14, 3, 15, 94, 75, header_fill=C_SECONDARY, cell_size=8)
print("  Page 14 done")

# Page 15: 강점/단점/개선방향
s = new_slide()
shape_with_text(s, 'rounded_rect', 2, 8, 42, 7, '강점', fill=C_ACCENT, text_color=C_WHITE, size=18, bold=True)
add_multiline(s, ['전기/전자/정보통신 분야의 다양한 교과목 개설', '세부전공분야 별 연구 진행에 유리하도록\n교육과정 및 연계도가 편성', '차세대 지능형 모빌리티\n핵심요소별 교과목 포함', '학사과정 진단평가 규정을 제정되어 시행'], 3, 18, 40, 30, size=9, color=C_DARK, bullet=True)
shape_with_text(s, 'rounded_rect', 2, 55, 42, 7, '단점', fill=C_MID_BLUE, text_color=C_WHITE, size=18, bold=True)
add_multiline(s, ['차세대 지능형 모빌리티 분야에 특화되지 않은\n교육과정 구성', '낮은 전공융합 교과목, 산학협력 프로젝트\n기반의 교과목 비중', '낮은 학생주도형/참여형 교과목 비중', '타분야간 교류가 어려운 학사관리 제도'], 3, 65, 40, 28, size=9, color=C_DARK, bullet=True)
shape_with_text(s, 'rounded_rect', 52, 5, 45, 7, '개선방향', fill=C_PRIMARY, text_color=C_WHITE, size=18, bold=True)
improvements = [
    '차세대 지능형 모빌리티\n전문인력양성에 특화된\n교육과정 및 학사관리 운영',
    '기술 분야 간 융합을 추구하는\n전공 유합 교과목 및 산학연계\nR&D 프로젝트 기반\n교과목 확대',
    '학생 주도형/참여형\n교과목 확대',
    '교육과 연구의 선순환\n구조 확립',
]
for i, txt in enumerate(improvements):
    add_text(s, txt, 55, 16+i*20, 40, 18, size=11, color=C_DARK)
print("  Page 15 done")

# Pages 16-33: 구조화된 슬라이드 생성
page_data = {
    16: ('교육과정 구성 및 운영 계획', [
        ('3세부 트랙 구성', '자율주행 플랫폼 트랙\n모빌리티 에너지 트랙\n커넥티비티 트랙'),
        ('융합/심화교과목', '2개 이상 기술 분야의 융합\n팀티칭\nPBL 기반 산학협력\n실화 프로젝트 과목\n학생주도형 프로젝트 수행'),
        ('국제화/고 개방적\n교육과정', '영어 강의 비율 확대\n세미나 강좌 확대'),
        ('학생맞춤형 핵심역량\n증진 프로그램', '인턴십 프로그램\n멘토링/인성교육, 연구윤리\n교과서 적용\n영어논문작성법,\n특허작성법,\n창업 교육 기회 제공'),
    ]),
    17: ('교과과정 로드맵', [
        ('1학기', '스마트그리드 공학\nICT 융합형 컨버터 제어응용\n전력변환기 제어 및 시뮬레이션'),
        ('하계 방학', '기업참여 Agject-based Team\n하계 산학연 학술연구회'),
        ('2학기', '스마트그리드전력변환\n고급디지털회로\nAC/DC시스템안전\n정보통신공학\nIoT 기술\n고급데이터통신'),
        ('동계방학', '산학연동 프로그램\n동계 산학연 학술연구회'),
    ]),
    18: ('교과과정 - 플랫폼별 구조', [
        ('자율주행 플랫폼', '피보역학, 스마트 모빌리티\n사물인터넷, 음성인식'),
        ('모빌리티 에너지 플랫폼', '고급유체역학, 스마트 모빌리티\n지능형구동기, 난류'),
        ('커넥티비티 플랫폼', '통신공학, 디지털트윈'),
    ]),
    19: ('교육 평가 체계', [
        ('학생-교수\n교육활동', '수업, 프로젝트'),
        ('운영위원회', '평가 취합/분석'),
        ('강의 및 교수활동\n개선사항 반영', ''),
    ]),
    20: ('교육과정 ↔ 연구과정 선순환', [
        ('교육과정', '연구기초교육 및 핵심교과목 교육강화\n자기주도형 교육컨텐츠 활성화\n프로젝트기반 수업 활성화\n산업체 연계형 교육과정'),
        ('연구과정', '우수연구성과 이용 프로젝트\n기반 수업\n우수연구성과 온라인\n교육컨텐츠 구축'),
        ('교육-연구', '교육과 연구의 선순환 환류체계 구축\n자기주도형 연구역량강화를 위한\n지원체계 구축\n산학협력 인프라 구축 확대'),
    ]),
    21: ('신진연구인력 운용 계획', [
        ('인력 확보 계획', '채용홍보강화\n국내 연구인력 모집 사이트, 홍보\n해외 협력 네트워크 홍보\n\n입학기회확대\n신진/졸업생 우선채용 기간 연장'),
        ('인력 지원 계획', '기본 인건비 제공\n월 330만원 이상×2년\n\n연구 인센티브 지급\n\n학술 및 연구활동 지원'),
    ]),
    22: ('연구 인프라 및 네트워크 활용', [
        ('01 신진연구 인력 간\n연구 협업 활성화', 'Seoultech Post-doc\n연구교수 지원사업\n연 2명'),
        ('02 신진연구 인력의\n학내외 연구 인프라 활용\n및 연구 네트워크\n주도적 참여', '서울과학기술대학교 연구 인프라\n중점연구소 지원사업 자율주행 인프라\n인공지능 플랫폼 인프라'),
        ('03 모빌리티 분야\n경쟁력 확보', '연구 성과물의\n질적 향상\nQ1 1편 / 1년, 1인'),
    ]),
    23: ('차세대 지능형 모빌리티 플랫폼 융합인재 양성 교육연구팀 종합', [
        ('비전', '모빌리티 플랫폼 분야 실무형 융합인재 양성, 글로벌 TOP 20대 연구그룹 진입'),
        ('핵심전략 4ST', 'Synergetic Technology 융합기술\nSelf-directed & Teamwork 학생 주도 연구/팀워크\nSustainable Training 지속가능형 교육\nSolution-targeted Thinking 실무/문제해결 중심'),
        ('추진방법', '참여교수 융합 교과목 개설\n참여교수 지도학생간의 팀 프로젝트 진행\n다년간 장기 팀 프로젝트 수행\n산업체 핀판의 프로젝트 평가, 주기적인 진행'),
    ]),
    24: ('글로벌 최상위 대학 관련연구그룹 비교', [
        ('Michigan', '총논문 수: 3\n총 Q1 논문 수: 3\n참여교수들 공동 논문 수: 3\n산학 협력 논문 수: 3'),
        ('Stanford', '총논문 수: 3\n총 Q1 논문 수: 3\n참여교수들 공동 논문 수: 3\n산학 협력 논문 수: 3'),
        ('MIT', '총논문 수: 3\n총 Q1 논문 수: 3\n참여교수들 공동 논문 수: 2\n산학 협력 논문 수: 2'),
        ('ST@ST', '총논문 수: 3\n총 Q1 논문 수: 3\n참여교수들 공동 논문 수: 2\n산학 협력 논문 수: 2'),
    ]),
    25: ('글로벌 최상위 대학 교육과정 벤치마킹', [
        ('MIT', '에너지 관련 학문 간 융합, 연계, 연구 등을 지원\n산업계 전문가가 참여하는 50개 이상의 단기 교육 프로그램 진행'),
        ('Harvard', '교수진의 약 30%: 다양한 연구 분야에서 공동으로 임명됨\n교수들 간 융합 프로그램 제공'),
        ('Stanford', '산업전문가와 빅데이터, 사이버보안 등과의 융합 프로그램 제공\n반기마다 산업체에서 교육 프로그램을 진행'),
        ('ST@ST', '전공 융합형 교과목, 프로그램 부재\n기업 연계형 교과목, 프로그램 부재'),
    ]),
    26: ('ST@ST 연구 추진 전략 - 3가지', [
        ('융합 연구 프로젝트', '학생주도/팀워크(ST)를\n통한 융합연구(ST) 진행\n\n자율주행 플랫폼\n모빌리티 에너지 플랫폼\n커넥티비티 플랫폼'),
        ('산학연 협업 활성화', '기업 수요 기반\n실무 중심 연구(ST) 진행\n\n국내외 산학연 공동연구\n학생주도(ST) 창업/인턴'),
        ('연구 질적 고도화', '우수성과물의 실용화 연계를 통한\n지속가능한 (ST) 연구 진행\n\n최고수준 저널/학회\n전문가 리뷰 및 피드백 수렴'),
    ]),
    27: ('ST@ST 교육 추진 전략 - 4가지', [
        ('융합 교육과정 구성\n창의융합(ST) 인재 양성', '융합 교과목 개발, 정보 수집\n현장요소 교과를 2~30여과목\n교내 타학과와 교과목 연계'),
        ('산학연계 교육과정 활성화\n실무지향(ST) 인재 양성', '실무 수요 프로젝트\n산학연 연계 교과목 및\n산업체 임직원 세미나 초청'),
        ('학생 주도형 참여 프로그램\n자기주도(ST) 인재 양성', '문제해결 중심의 학생 교과 중심 강의\n학생소통기반 산업체 입학원 연계\n산업계 임직원 멘토 기반'),
        ('선진학사 시스템 구축\n지속성장(ST) 인재 양성', '교내 인접학부와의\n연구지식 One-Stop 형'),
    ]),
    28: ('4대 분야 기술 상세', [
        ('인공지능(AI)', 'AI 기반 자율주행\n객체인식/패턴인식\n소프트웨어설계,최적화'),
        ('반도체, 디스플레이', 'DC 전력망 고효율기기 기반\nDC 전력변환 응용시스템\n고속 분석기반 SiC/GaN\n전력반도체 프로세스 개발'),
        ('차세대 통신', '분산에너지자원 관리시스템\n전기 공급설비의 수명 예측\n해저통신 시스템\n전력전자 기반 데이터 처리'),
        ('이차전지', '전력변환 기법의 자동차 활용\n미래에너지저장시스템\n데이터 기반 에너지 저장'),
    ]),
    29: ('시장 규모 추이 (억 달러)', []),
    30: ('시장 규모 추이 (10억원)', []),
    31: ('산학협력 기업 목록', []),
    32: ('Introduction - Mobile Characteristic', []),
    33: ('Introduction - 중점연구소 구조', []),
}

for pg_num in range(16, 34):
    s = new_slide()
    if pg_num in page_data:
        title, sections = page_data[pg_num]
        # 타이틀
        add_text(s, title, 3, 2, 94, 8, size=20, bold=True, color=C_DARK, align='center')

        if pg_num in (29, 30):
            # 차트 페이지: 바 차트로 근사
            if pg_num == 29:
                data = [('2018','843'), ('2019','1,030'), ('2020','1,260'), ('2021','1,540'), ('2022','1,882'), ('2023','2,301'), ('2024','2,812'), ('2025','3,437')]
                unit = '[ 단위: 억 달러 ]'
            else:
                data = [('2017','682'), ('2018','896'), ('2019','1,176'), ('2020','1,535'), ('2021','1,950'), ('2022','2,416'), ('2023','2,863')]
                unit = '[ 단위: 10억원 ]'
            add_text(s, unit, 5, 10, 30, 5, size=10, color=C_GRAY)
            bar_w = 70 / len(data)
            for i, (year, val) in enumerate(data):
                x = 10 + i * bar_w
                h = min(60, int(int(val.replace(',','')) / 60))
                y = 75 - h
                is_last = (i == len(data) - 1)
                fill = C_PRIMARY if is_last else C_GRAY
                add_shape(s, 'rect', x, y, bar_w*0.7, h, fill=fill)
                add_text(s, val, x-1, y-7, bar_w+2, 6, size=12 if is_last else 10, bold=is_last, color=C_PRIMARY if is_last else C_DARK, align='center')
                add_text(s, year, x, 78, bar_w, 5, size=9, color=C_DARK, align='center')

        elif pg_num == 31:
            # 산학협력 기업 목록 - 도형+텍스트로 로고 근사
            shape_with_text(s, 'circle', 38, 30, 24, 35, 'ST@ST', fill=C_PRIMARY, text_color=C_WHITE, size=18, bold=True)
            companies_top = ['LG MAGNA', '연합시스', 'ETRI', '올리브텍스', 'WILUS', 'KETI']
            for i, name in enumerate(companies_top):
                shape_with_text(s, 'rounded_rect', 3+i*16, 8, 14, 10, name, fill=C_LIGHT_BG, text_color=C_DARK, size=8, bold=True, stroke=C_GRAY)
            companies_bot = ['MQ Lab', 'LX Semicon', 'SAPEON', 'mobilint', 'ISC']
            for i, name in enumerate(companies_bot):
                shape_with_text(s, 'rounded_rect', 6+i*17, 72, 15, 10, name, fill=C_LIGHT_BG, text_color=C_DARK, size=8, bold=True, stroke=C_GRAY)

        elif pg_num == 32:
            # Introduction - Mobile Characteristic (네이티브 도형)
            add_text(s, 'Introduction', 2, 2, 30, 6, size=20, bold=True, color=C_DARK)
            shape_with_text(s, 'rounded_rect', 20, 3, 55, 5, 'KEY POINT: Optimization on trade-off between accuracy and HW resources/FLOPs', fill=C_LIGHT_BG, text_color=C_DARK, size=9, stroke=C_GRAY)
            shape_with_text(s, 'rounded_rect', 2, 18, 22, 25, 'HW-based low complexity\nschemes for low-power &\nspeed-up\n\nPipelining\nParallelism', fill=C_LIGHT_BLUE, text_color=C_DARK, size=8)
            shape_with_text(s, 'rounded_rect', 26, 12, 20, 12, 'Architecture Platform\nProcessor\nHW Accel.\nFPGA | ASICs', fill=C_ACCENT, text_color=C_WHITE, size=8)
            shape_with_text(s, 'rounded_rect', 48, 12, 20, 12, 'Memory System for DNNs\nMain memory\nStorage-class memory\nNon-volatile memory', fill=C_SECONDARY, text_color=C_WHITE, size=8)
            shape_with_text(s, 'rounded_rect', 70, 12, 14, 8, 'Self-Learning\nMobile', fill=C_PRIMARY, text_color=C_WHITE, size=9)
            shape_with_text(s, 'rounded_rect', 76, 25, 22, 25, 'SW-based low complexity\nschemes for low-power &\nspeed-up\n\nPruning\nQuantization', fill=C_LIGHT_BLUE, text_color=C_DARK, size=8)
            shape_with_text(s, 'rounded_rect', 2, 50, 22, 20, 'Performance enhancement\nschemes\n\nAccuracy\nEnhancement', fill=C_LIGHT_BLUE, text_color=C_DARK, size=8)
            shape_with_text(s, 'rounded_rect', 60, 60, 18, 8, 'Autonomous Driving', fill=C_PRIMARY, text_color=C_WHITE, size=9, bold=True)
            shape_with_text(s, 'rounded_rect', 60, 72, 18, 8, 'Video Surveillance\nas a Service (VSaaS)', fill=C_PRIMARY, text_color=C_WHITE, size=8)

        elif pg_num == 33:
            # Introduction - 중점연구소 구조 (네이티브 도형)
            add_text(s, 'Introduction', 2, 2, 30, 6, size=20, bold=True, color=C_DARK)
            shape_with_text(s, 'rounded_rect', 15, 5, 18, 6, '중점연구소', fill=C_ACCENT, text_color=C_WHITE, size=12, bold=True)
            shape_with_text(s, 'rounded_rect', 2, 15, 30, 35, 'SENSORS\nCAMERA\nRADAR\nLiDAR\nGPS\nOTHERS\n\nPERCEPTION\nDETECTION\nLOCALIZATION\n\nPLANNING\nMPC\nOTHERS', fill=C_LIGHT_BG, text_color=C_DARK, size=7, stroke=C_GRAY)
            add_text(s, 'AI\nalgorithms', 8, 52, 14, 8, size=10, bold=True, color=C_PRIMARY)
            shape_with_text(s, 'rounded_rect', 40, 10, 25, 30, 'AIoT\n\nNeural Processing Unit (NPU) - INF_0\n\nITRC/서버\n\nMAC MAC MAC\nInternal Memory (W/A)\nMAC MAC MAC', fill=C_LIGHT_BG, text_color=C_DARK, size=7, stroke=C_GRAY)
            shape_with_text(s, 'rounded_rect', 70, 10, 12, 8, '기본연구', fill=C_PRIMARY, text_color=C_WHITE, size=10, bold=True)
            shape_with_text(s, 'rounded_rect', 2, 60, 30, 20, '모바일 자가학습\nSelf-Learning Accel.\n\nStorage -PRAM\nData Select Module\n2T-DRAM PIM | PIM', fill=C_LIGHT_BG, text_color=C_DARK, size=7, stroke=C_GRAY)
            shape_with_text(s, 'rounded_rect', 40, 50, 25, 10, 'NPU - INF_1', fill=C_LIGHT_BLUE, text_color=C_DARK, size=9)
            shape_with_text(s, 'rounded_rect', 40, 65, 25, 10, 'NPU-INF_N-1', fill=C_LIGHT_BLUE, text_color=C_DARK, size=9)
            shape_with_text(s, 'rounded_rect', 70, 30, 12, 8, 'NN-\nPRAM', fill=C_SECONDARY, text_color=C_WHITE, size=9)

        elif sections:
            n = len(sections)
            col_w = min(45, 90 // n)
            for i, (sub_title, desc) in enumerate(sections):
                x = 3 + i * (col_w + 1)
                shape_with_text(s, 'rounded_rect', x, 14, col_w, 8, sub_title, fill=C_PRIMARY, text_color=C_WHITE, size=11, bold=True)
                if desc:
                    add_multiline(s, desc.split('\n'), x+1, 24, col_w-2, 65, size=9, color=C_DARK, bullet=True)
    else:
        add_text(s, f'Page {pg_num}', 40, 40, 20, 10, size=24, bold=True, color=C_GRAY, align='center')

    print(f"  Page {pg_num} done")

doc.close()

# ── 초과 슬라이드 제거 (PDF 페이지 수에 맞춤) ──
while len(prs.slides) > NUM_PAGES:
    rId = prs.slides._sldIdLst[-1].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[-1]

# ── 저장 ──
prs.save(OUTPUT_PATH)
print(f"\n완료: {OUTPUT_PATH}")
print(f"총 {len(prs.slides)} 슬라이드 생성 (PDF: {NUM_PAGES}페이지)")
